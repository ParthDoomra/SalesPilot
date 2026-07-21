import os
import io
import logging
from typing import Dict, Any, List
from pypdf import PdfReader
import docx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from fpdf import FPDF

from app.database import bucket
from app.config import settings

logger = logging.getLogger("salespilot.document")

class DocumentService:
    def __init__(self):
        # Ensure local exports folder exists for mock downloads
        self.local_export_dir = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", "static", "exports"))
        os.makedirs(self.local_export_dir, exist_ok=True)

    def extract_text_from_file(self, file_content: bytes, file_name: str) -> str:
        """
        Parses PDF, DOCX, or TXT content and extracts full string contents.
        """
        ext = file_name.split(".")[-1].lower()
        
        if ext == "pdf":
            try:
                pdf_file = io.BytesIO(file_content)
                reader = PdfReader(pdf_file)
                text = ""
                for page in reader.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text + "\n"
                return text
            except Exception as e:
                logger.error(f"Error parsing PDF: {e}")
                raise ValueError("Failed to parse PDF document.")
                
        elif ext in ["docx", "doc"]:
            try:
                docx_file = io.BytesIO(file_content)
                doc = docx.Document(docx_file)
                text = ""
                for para in doc.paragraphs:
                    text += para.text + "\n"
                return text
            except Exception as e:
                logger.error(f"Error parsing DOCX: {e}")
                raise ValueError("Failed to parse DOCX document.")
                
        elif ext == "txt":
            try:
                return file_content.decode("utf-8", errors="ignore")
            except Exception as e:
                logger.error(f"Error parsing TXT: {e}")
                raise ValueError("Failed to parse TXT document.")
                
        else:
            raise ValueError(f"Unsupported file format: {ext}")

    def generate_docx_proposal(self, project: Dict[str, Any], content: Dict[str, Any]) -> str:
        """
        Generates a DOCX file containing proposal details.
        Returns the download URL (local static URL if demo mode, Firebase Storage URL if prod).
        """
        doc = docx.Document()
        
        # Style / Document Title
        doc.add_heading(f"SalesPilot AI Enterprise Solution Proposal", 0)
        doc.add_paragraph(f"Client: {project.get('clientName', 'N/A')} ({project.get('company', 'N/A')})")
        doc.add_paragraph(f"Industry: {project.get('industry', 'N/A')} | Country: {project.get('country', 'N/A')}")
        doc.add_paragraph(f"Budget Limit: ${project.get('budget', 0):,}")
        doc.add_paragraph(f"Timeline: {project.get('timeline', 'N/A')}")
        
        # Sections
        # 1. Executive Summary
        doc.add_heading("1. Executive Summary", level=1)
        doc.add_paragraph(content.get("executiveSummary", "This proposal outlines the architected solution designed for cloud deployment."))

        # 2. Requirements Summary
        doc.add_heading("2. Technical & Business Requirements", level=1)
        reqs = content.get("requirements", {})
        for req_cat, req_list in reqs.items():
            if isinstance(req_list, list):
                doc.add_heading(req_cat.capitalize(), level=2)
                for item in req_list:
                    doc.add_paragraph(item, style='List Bullet')
            else:
                doc.add_paragraph(f"{req_cat}: {req_list}")

        # 3. Cloud Architecture
        doc.add_heading("3. Cloud Architecture & Tech Stack", level=1)
        arch = content.get("architecture", {})
        doc.add_paragraph(f"Preferred Cloud: {project.get('preferredCloud', 'Azure')}")
        for svc_cat, svc_detail in arch.get("techStack", {}).items():
            doc.add_heading(svc_cat.capitalize(), level=2)
            doc.add_paragraph(f"Service: {svc_detail.get('service', 'N/A')}")
            doc.add_paragraph(f"Rationale: {svc_detail.get('rationale', 'N/A')}")

        # 4. Pricing & Bill of Materials
        doc.add_heading("4. Financial Estimate (Bill of Materials)", level=1)
        pricing = content.get("pricing", {})
        doc.add_paragraph(f"Total Estimated Monthly Cost: ${pricing.get('monthlyTotal', 0.0):,}")
        doc.add_paragraph(f"Total Estimated Annual Cost: ${pricing.get('annualTotal', 0.0):,}")
        
        # Table of costs
        table = doc.add_table(rows=1, cols=2)
        hdr_cells = table.rows[0].cells
        hdr_cells[0].text = 'Service Category'
        hdr_cells[1].text = 'Monthly Cost (USD)'
        
        for cat, val in pricing.get("breakdown", {}).items():
            row_cells = table.add_row().cells
            row_cells[0].text = cat.capitalize()
            row_cells[1].text = f"${val:,}"

        # 5. Resilience & Recovery Plan
        doc.add_heading("5. Resilience & Failure Recovery Report", level=1)
        resilience = content.get("resilience", {})
        doc.add_paragraph(f"Resilience Rating Score: {resilience.get('score', 0)} / 100")
        doc.add_paragraph("Disaster Scenarios Simulated:")
        for sc in resilience.get("scenarios", []):
            p = doc.add_paragraph()
            p.add_run(f"- {sc.get('name')}: ").bold = True
            p.add_run(f"Downtime (RTO): {sc.get('rto')}, Data Loss (RPO): {sc.get('rpo')}. Impact: {sc.get('impact')}")

        # Save file
        filename = f"proposal_{project.get('id', 'new')}.docx"
        local_path = os.path.join(self.local_export_dir, filename)
        doc.save(local_path)

        return self._get_export_url(filename, local_path, "application/vnd.openxmlformats-officedocument.wordprocessingml.document")

    def generate_pdf_proposal(self, project: Dict[str, Any], content: Dict[str, Any]) -> str:
        """
        Generates a PDF file containing proposal details.
        """
        pdf = FPDF()
        pdf.add_page()
        pdf.set_auto_page_break(auto=True, margin=15)
        
        # Title
        pdf.set_font("helvetica", "B", 18)
        pdf.cell(0, 10, "SalesPilot AI - Cloud Architecture Proposal", ln=True, align="C")
        pdf.set_font("helvetica", "", 10)
        pdf.cell(0, 8, f"Client: {project.get('clientName', 'N/A')} ({project.get('company', 'N/A')})", ln=True, align="C")
        pdf.ln(5)

        # Draw a horizontal line
        pdf.line(10, pdf.get_y(), 200, pdf.get_y())
        pdf.ln(5)

        # 1. Executive Summary
        pdf.set_font("helvetica", "B", 14)
        pdf.cell(0, 10, "1. Executive Summary", ln=True)
        pdf.set_font("helvetica", "", 10)
        summary_text = content.get("executiveSummary", "Comprehensive cloud architecture solution.")
        pdf.multi_cell(0, 6, summary_text)
        pdf.ln(5)

        # 2. Project Metadata
        pdf.set_font("helvetica", "B", 14)
        pdf.cell(0, 10, "2. Project Specifications", ln=True)
        pdf.set_font("helvetica", "", 10)
        pdf.cell(0, 6, f"Preferred Cloud Provider: {project.get('preferredCloud', 'Azure')}", ln=True)
        pdf.cell(0, 6, f"Client Budget Cap: ${project.get('budget', 0):,} USD", ln=True)
        pdf.cell(0, 6, f"Project Timeline: {project.get('timeline', 'N/A')}", ln=True)
        pdf.ln(5)

        # 3. Pricing
        pdf.set_font("helvetica", "B", 14)
        pdf.cell(0, 10, "3. Cost Estimates (Monthly)", ln=True)
        pdf.set_font("helvetica", "", 10)
        pricing = content.get("pricing", {})
        pdf.cell(0, 6, f"Total Monthly Cost: ${pricing.get('monthlyTotal', 0.0):,} USD", ln=True)
        pdf.cell(0, 6, f"Total Annualized Cost: ${pricing.get('annualTotal', 0.0):,} USD", ln=True)
        pdf.ln(2)
        
        # Draw small pricing table
        pdf.set_font("helvetica", "B", 10)
        pdf.cell(80, 8, "Category", border=1)
        pdf.cell(60, 8, "Estimated Cost (USD)", border=1, ln=True)
        pdf.set_font("helvetica", "", 10)
        for cat, val in pricing.get("breakdown", {}).items():
            pdf.cell(80, 8, cat.capitalize(), border=1)
            pdf.cell(60, 8, f"${val:,}", border=1, ln=True)
        pdf.ln(5)

        # 4. Resilience
        pdf.set_font("helvetica", "B", 14)
        pdf.cell(0, 10, "4. Solution Architecture Resilience", ln=True)
        pdf.set_font("helvetica", "", 10)
        resilience = content.get("resilience", {})
        pdf.cell(0, 6, f"Infrastructure Resilience Score: {resilience.get('score', 0)} / 100", ln=True)
        pdf.ln(2)

        # Save
        filename = f"proposal_{project.get('id', 'new')}.pdf"
        local_path = os.path.join(self.local_export_dir, filename)
        pdf.output(local_path)

        return self._get_export_url(filename, local_path, "application/pdf")

    def generate_pptx_slides(self, project: Dict[str, Any], content: Dict[str, Any]) -> str:
        """
        Generates an editable PPTX deck with slides:
        - Title / Client Overview
        - Requirements Summary
        - Cloud Architecture
        - Pricing Breakdown
        - resilience & Disaster Recovery
        - Optimizations / Savings (Negotiation)
        - Conclusion
        """
        prs = Presentation()
        
        # Color Palette
        DARK_BG = RGBColor(15, 23, 42)     # Tailwind Slate-900
        WHITE = RGBColor(255, 255, 255)
        GREY_TEXT = RGBColor(148, 163, 184) # Tailwind Slate-400
        CYAN = RGBColor(6, 182, 212)       # Tailwind Cyan-500
        
        # Helper to set background color and add titles
        def apply_dark_theme(slide):
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = DARK_BG

        def add_slide_title(slide, text: str):
            tx_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9.0), Inches(1.0))
            tf = tx_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = text
            p.font.name = 'Helvetica'
            p.font.size = Pt(36)
            p.font.bold = True
            p.font.color.rgb = CYAN

        # --- Slide 1: Title ---
        slide_layout = prs.slide_layouts[6] # Blank
        slide = prs.slides.add_slide(slide_layout)
        apply_dark_theme(slide)
        
        # Main Title & Subtitle
        tx_box = slide.shapes.add_textbox(Inches(0.75), Inches(2.0), Inches(8.5), Inches(3.0))
        tf = tx_box.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = f"SalesPilot AI Solution Proposal"
        p.font.name = 'Helvetica'
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = WHITE
        
        p2 = tf.add_paragraph()
        p2.text = f"Enterprise Cloud Architecture for {project.get('company', 'Acme Co')}"
        p2.font.name = 'Helvetica'
        p2.font.size = Pt(22)
        p2.font.color.rgb = CYAN
        p2.space_before = Pt(20)

        p3 = tf.add_paragraph()
        p3.text = f"Prepared for: {project.get('clientName', 'Executive Team')} | Date: 2026"
        p3.font.name = 'Helvetica'
        p3.font.size = Pt(14)
        p3.font.color.rgb = GREY_TEXT
        p3.space_before = Pt(40)

        # --- Slide 2: Client & Requirements ---
        slide = prs.slides.add_slide(slide_layout)
        apply_dark_theme(slide)
        add_slide_title(slide, "Technical & Business Requirements")
        
        tx_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(9.0), Inches(4.5))
        tf = tx_box.text_frame
        tf.word_wrap = True
        
        reqs = content.get("requirements", {})
        # Flatten requirements items
        req_lines = []
        if isinstance(reqs, dict):
            for cat, val in reqs.items():
                if isinstance(val, list):
                    req_lines.extend(val[:2])
                else:
                    req_lines.append(f"{cat}: {val}")
        
        for line in req_lines[:6]:
            p = tf.add_paragraph() if tf.text else tf.paragraphs[0]
            p.text = f"• {line}"
            p.font.name = 'Helvetica'
            p.font.size = Pt(18)
            p.font.color.rgb = WHITE
            p.space_after = Pt(12)

        # --- Slide 3: Proposed Tech Stack ---
        slide = prs.slides.add_slide(slide_layout)
        apply_dark_theme(slide)
        add_slide_title(slide, "Proposed Solution Architecture")
        
        tx_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(9.0), Inches(4.5))
        tf = tx_box.text_frame
        tf.word_wrap = True
        
        arch = content.get("architecture", {})
        tech_stack = arch.get("techStack", {})
        
        count = 0
        for category, service_info in tech_stack.items():
            if count >= 4:
                break
            p = tf.add_paragraph() if tf.text else tf.paragraphs[0]
            p.text = f"{category.capitalize()}: {service_info.get('service', 'N/A')}"
            p.font.name = 'Helvetica'
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = CYAN
            p.space_before = Pt(8)
            
            p_sub = tf.add_paragraph()
            p_sub.text = f"   Rationale: {service_info.get('rationale', 'N/A')}"
            p_sub.font.name = 'Helvetica'
            p_sub.font.size = Pt(14)
            p_sub.font.color.rgb = WHITE
            p_sub.space_after = Pt(12)
            count += 1

        # --- Slide 4: Cost Breakdown ---
        slide = prs.slides.add_slide(slide_layout)
        apply_dark_theme(slide)
        add_slide_title(slide, "Financial Estimate & BOM")
        
        pricing = content.get("pricing", {})
        
        # Left side: Total metrics
        tx_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(4.0), Inches(4.0))
        tf_left = tx_box.text_frame
        tf_left.word_wrap = True
        
        p = tf_left.paragraphs[0]
        p.text = "Monthly Investment"
        p.font.name = 'Helvetica'
        p.font.size = Pt(18)
        p.font.color.rgb = GREY_TEXT
        
        p_val = tf_left.add_paragraph()
        p_val.text = f"${pricing.get('monthlyTotal', 0.0):,}"
        p_val.font.name = 'Helvetica'
        p_val.font.size = Pt(40)
        p_val.font.bold = True
        p_val.font.color.rgb = WHITE
        p_val.space_after = Pt(24)
        
        p_ann = tf_left.add_paragraph()
        p_ann.text = "Annualized Total"
        p_ann.font.name = 'Helvetica'
        p_ann.font.size = Pt(18)
        p_ann.font.color.rgb = GREY_TEXT
        
        p_ann_val = tf_left.add_paragraph()
        p_ann_val.text = f"${pricing.get('annualTotal', 0.0):,}"
        p_ann_val.font.name = 'Helvetica'
        p_ann_val.font.size = Pt(32)
        p_ann_val.font.bold = True
        p_ann_val.font.color.rgb = CYAN

        # Right side: Detail List
        tx_box_right = slide.shapes.add_textbox(Inches(5.0), Inches(2.0), Inches(4.5), Inches(4.0))
        tf_right = tx_box_right.text_frame
        tf_right.word_wrap = True
        
        for item_cat, item_val in pricing.get("breakdown", {}).items():
            p_item = tf_right.add_paragraph() if tf_right.text else tf_right.paragraphs[0]
            p_item.text = f"• {item_cat.capitalize()}: ${item_val:,} / mo"
            p_item.font.name = 'Helvetica'
            p_item.font.size = Pt(18)
            p_item.font.color.rgb = WHITE
            p_item.space_after = Pt(10)

        # --- Slide 5: Resilience & Risk ---
        slide = prs.slides.add_slide(slide_layout)
        apply_dark_theme(slide)
        add_slide_title(slide, "Infrastructure Resilience & Recovery")
        
        resilience = content.get("resilience", {})
        
        tx_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(9.0), Inches(4.5))
        tf = tx_box.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = f"Resilience Score: {resilience.get('score', 0)} / 100"
        p.font.name = 'Helvetica'
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = CYAN
        p.space_after = Pt(20)
        
        for sc in resilience.get("scenarios", [])[:3]:
            p_sc = tf.add_paragraph()
            p_sc.text = f"Scenario: {sc.get('name')}"
            p_sc.font.name = 'Helvetica'
            p_sc.font.size = Pt(18)
            p_sc.font.bold = True
            p_sc.font.color.rgb = WHITE
            
            p_sc_det = tf.add_paragraph()
            p_sc_det.text = f"   Impact: {sc.get('impact')} | RTO: {sc.get('rto')} | RPO: {sc.get('rpo')}"
            p_sc_det.font.name = 'Helvetica'
            p_sc_det.font.size = Pt(14)
            p_sc_det.font.color.rgb = GREY_TEXT
            p_sc_det.space_after = Pt(10)

        # --- Slide 6: Optimizations (Negotiation) ---
        slide = prs.slides.add_slide(slide_layout)
        apply_dark_theme(slide)
        add_slide_title(slide, "Architecture Optimization Suggestions")
        
        neg = content.get("negotiation", {})
        
        tx_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(9.0), Inches(4.5))
        tf = tx_box.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = f"Original Cost: ${neg.get('originalCost', pricing.get('monthlyTotal', 0.0)):,} | Optimized Cost: ${neg.get('optimizedCost', 0.0):,}"
        p.font.name = 'Helvetica'
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = CYAN
        p.space_after = Pt(16)
        
        for mod in neg.get("modifications", [])[:3]:
            p_mod = tf.add_paragraph()
            p_mod.text = f"• Optimization: {mod.get('suggestion')}"
            p_mod.font.name = 'Helvetica'
            p_mod.font.size = Pt(18)
            p_mod.font.color.rgb = WHITE
            
            p_trade = tf.add_paragraph()
            p_trade.text = f"  Savings: {mod.get('savings')} | Trade-offs: {mod.get('tradeOff')}"
            p_trade.font.name = 'Helvetica'
            p_trade.font.size = Pt(14)
            p_trade.font.color.rgb = GREY_TEXT
            p_trade.space_after = Pt(8)

        # Save PowerPoint
        filename = f"proposal_{project.get('id', 'new')}.pptx"
        local_path = os.path.join(self.local_export_dir, filename)
        prs.save(local_path)
        
        return self._get_export_url(filename, local_path, "application/vnd.openxmlformats-officedocument.presentationml.presentation")

    def _get_export_url(self, filename: str, local_path: str, content_type: str) -> str:
        """
        Resolves a downloadable URL for the generated files.
        Uploads to Firebase Storage if configured and not in demo mode.
        """
        if settings.is_firebase_configured and not settings.DEMO_MODE:
            try:
                blob_name = f"proposals/{filename}"
                blob = bucket.blob(blob_name)
                blob.upload_from_filename(local_path, content_type=content_type)
                blob.make_public()
                return blob.public_url
            except Exception as e:
                logger.error(f"Failed to upload document to Firebase Storage: {e}. Serving locally.")
        
        # fallback: serve locally via FastAPI static mount
        return f"/static/exports/{filename}"

document_service = DocumentService()
